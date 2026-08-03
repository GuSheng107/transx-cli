from pathlib import Path
import sys

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from reportlab.pdfgen.canvas import Canvas


output = Path(sys.argv[1])
output.mkdir(parents=True, exist_ok=True)

document = Document()
paragraph = document.add_paragraph()
bold = paragraph.add_run("Research title. ")
bold.bold = True
normal = paragraph.add_run("This document contains two formatted text runs.")
normal.font.name = "Arial"
table = document.add_table(rows=1, cols=2)
table.cell(0, 0).text = "Method"
table.cell(0, 1).text = "Result"
document.save(output / "paper.docx")

workbook = Workbook()
sheet = workbook.active
sheet.title = "Research"
sheet["A1"] = "Sample"
sheet["B1"] = "Description"
sheet["A1"].font = Font(bold=True)
sheet["A1"].fill = PatternFill("solid", fgColor="FFF2CC")
sheet["A2"] = "A"
sheet["B2"] = "Short research result"
sheet["C2"] = "=1+1"
workbook.save(output / "data.xlsx")

presentation = Presentation()
slide = presentation.slides.add_slide(presentation.slide_layouts[6])
box = slide.shapes.add_textbox(914400, 914400, 7315200, 1828800)
paragraph = box.text_frame.paragraphs[0]
run = paragraph.add_run()
run.text = "Research presentation"
run.font.bold = True
run.font.size = Pt(24)
run.font.color.rgb = RGBColor(31, 78, 121)
run = paragraph.add_run()
run.text = " with formatted content"
run.font.size = Pt(18)
presentation.save(output / "slides.pptx")

pdf = Canvas(str(output / "paper.pdf"))
pdf.drawString(72, 720, "This PDF contains a short research abstract.")
pdf.save()

(output / "notes.txt").write_text("First research note.\nSecond research note.\n", encoding="utf-8")
(output / "guide.md").write_text("# Research guide\n\nRead [source](https://example.com).\n\n```js\nconst keep = true;\n```\n", encoding="utf-8")
(output / "data.csv").write_text('name,note\nA,"Short, clear result"\n', encoding="utf-8")
